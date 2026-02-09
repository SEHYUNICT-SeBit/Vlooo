"""
PPT 파싱 서비스
python-pptx를 사용한 PowerPoint 파일 분석
"""

import io
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
from PIL import Image


class Slide:
    """슬라이드 데이터 모델"""
    def __init__(self, slide_number: int, title: Optional[str] = None):
        self.slide_id = f"slide_{slide_number}"
        self.slide_number = slide_number
        self.title = title or f"Slide {slide_number}"
        self.content = ""
        self.image_data: List[bytes] = []
        self.notes = ""
        
    def to_dict(self) -> Dict[str, Any]:
        return {
            "slideId": self.slide_id,
            "slideNumber": self.slide_number,
            "title": self.title,
            "content": self.content,
            "imageUrls": [],  # 나중에 R2에 업로드 후 채움
            "notes": self.notes,
        }


class PptParser:
    """PowerPoint 파일 파서"""
    
    def __init__(self):
        self.prs: Optional[Presentation] = None
        self.slides: List[Slide] = []
    
    def parse_file(self, file_content: bytes) -> Dict[str, Any]:
        """PPT 파일 파싱"""
        try:
            # 바이너리 데이터로부터 Presentation 로드
            prs = Presentation(io.BytesIO(file_content))
            self.prs = prs
            
            # 슬라이드 추출
            self._extract_slides()
            
            return {
                "success": True,
                "total_slides": len(self.slides),
                "slides": [slide.to_dict() for slide in self.slides],
                "extracted_text": self._get_all_text(),
                "metadata": self._extract_metadata(),
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }
    
    def _extract_slides(self) -> None:
        """모든 슬라이드 추출"""
        if not self.prs:
            return
        
        for idx, slide in enumerate(self.prs.slides, 1):
            slide_obj = Slide(idx)
            
            # 슬라이드 제목 추출
            slide_obj.title = self._extract_title(slide)
            
            # 텍스트 추출
            slide_obj.content = self._extract_text(slide)
            
            # 발표자 노트 추출
            if slide.has_notes_slide:
                slide_obj.notes = self._extract_notes(slide)
            
            # 이미지 추출
            slide_obj.image_data = self._extract_images(slide)
            
            self.slides.append(slide_obj)
    
    def _extract_title(self, slide) -> str:
        """슬라이드 제목 추출"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except:
            pass
        return ""
    
    def _extract_text(self, slide) -> str:
        """슬라이드의 모든 텍스트 추출"""
        text_parts = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # 제목은 이미 추출했으므로 스킵
                    if shape != slide.shapes.title:
                        text_parts.append(shape.text.strip())
        except:
            pass
        
        return "\n".join(text_parts)
    
    def _extract_notes(self, slide) -> str:
        """발표자 노트 추출"""
        try:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        except:
            pass
        return ""
    
    def _extract_images(self, slide) -> List[bytes]:
        """슬라이드에서 이미지 추출"""
        images = []
        
        try:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    try:
                        image_stream = shape.image.blob
                        images.append(image_stream)
                    except:
                        pass
        except:
            pass
        
        return images
    
    def _get_all_text(self) -> str:
        """모든 슬라이드의 텍스트 통합"""
        all_text = []
        for slide in self.slides:
            if slide.title:
                all_text.append(f"# {slide.title}")
            if slide.content:
                all_text.append(slide.content)
            all_text.append("")
        
        return "\n".join(all_text)
    
    def _extract_metadata(self) -> Dict[str, Any]:
        """파일 메타데이터 추출"""
        if not self.prs:
            return {}
        
        try:
            core_props = self.prs.core_properties
            return {
                "pptTitle": core_props.title or "Untitled",
                "pptAuthor": core_props.author or "Unknown",
                "createdAt": str(core_props.created or ""),
            }
        except:
            return {}


def parse_ppt_file(file_content: bytes) -> Dict[str, Any]:
    """PPT 파일을 파싱하여 반환"""
    parser = PptParser()
    return parser.parse_file(file_content)
